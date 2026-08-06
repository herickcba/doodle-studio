#!/usr/bin/env python3
"""build_deck.py -- gera o deck de documentacao do design system.

    python3 tools/design-system/build_deck.py [saida.pptx]

Depois de gerar, VALIDA o proprio arquivo: le cada run de volta e confere fonte,
corpo, cor e entrelinha contra os tokens. Mesmo principio do vba-static-scan --
o gerador nao pode emitir nada fora do sistema, e a checagem e' no artefato
final, nao na intencao do codigo.
"""

from __future__ import annotations

import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.abspath(os.path.join(HERE, "..", ".."))
sys.path.insert(0, HERE)

from pptx import Presentation  # noqa: E402
from pptx.util import Pt  # noqa: E402

import layouts  # noqa: E402
from layouts import CONTENT_END, FOOTER_H  # noqa: E402
import spec as spec_mod  # noqa: E402
import tokens as T  # noqa: E402

FOOTER_Y = T.CONTENT_BOTTOM_PT - FOOTER_H

OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
    REPO, "docs", "CBA-Studio-Design-System-2.2.pptx")


def generate(path):
    prs = Presentation()
    prs.slide_width = Pt(T.PAGE_W_PT)
    prs.slide_height = Pt(T.PAGE_H_PT)

    slides = spec_mod.build_spec()
    for i, s in enumerate(slides, 1):
        try:
            layouts.build(prs, s, i)
        except Exception as e:
            raise SystemExit("ERRO no slide %d (%s): %s" % (i, s.get("kind"), e))

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    return len(slides)


def validate(path):
    """Le o PPTX de volta e confere cada run contra os tokens.

    Alem dos tokens, confere TRANSBORDO: mede o texto com a fonte real e
    compara com a caixa. Texto que nao cabe sobrepoe o bloco de baixo -- foi o
    unico defeito que a validacao de tokens deixou passar na primeira versao.
    """
    import measure  # noqa: E402

    prs = Presentation(path)
    sizes_ok = {float(s) for s in T.ALLOWED_SIZES}
    hexes_ok = {h.upper() for h in T.ALLOWED_HEX}
    ents_ok = {round(s["ent"], 3) for s in T.STYLES}
    by_size = {}
    for st in T.STYLES:
        by_size.setdefault((st["size"], bool(st["bold"]), round(st["ent"], 3)), st)

    errs, warns = [], []
    runs = 0
    for n, slide in enumerate(prs.slides, 1):
        content_bottom = 0.0
        try:
            slide_bg = str(slide.background.fill.fore_color.rgb).upper()
        except Exception:
            slide_bg = None
        for shape in slide.shapes:
            # ---- card da cor do fundo nao e' card. So' vale com contorno,
            # que e' o caso da amostra de cor da propria cor do fundo.
            if slide_bg and not (shape.has_text_frame
                                 and shape.text_frame.text.strip()):
                try:
                    fill_hex = str(shape.fill.fore_color.rgb).upper()
                except Exception:
                    fill_hex = None
                if fill_hex == slide_bg:
                    outlined = False
                    try:
                        outlined = shape.line.color.rgb is not None
                    except Exception:
                        pass
                    if not outlined:
                        errs.append("slide %d: forma da mesma cor do fundo (%s) "
                                    "e sem contorno" % (n, fill_hex))
            # ---- caixa esticada: forma com muito mais altura que conteudo.
            # Era o vicio de puxar o card ate' o rodape (item 7 da revisao).
            if (not shape.has_text_frame or not shape.text_frame.text.strip()) \
                    and shape.height and shape.width:
                bh = shape.height / 12700
                tallest = 0.0
                for other in slide.shapes:
                    if other is shape or not other.has_text_frame:
                        continue
                    if not other.text_frame.text.strip():
                        continue
                    ox, oy = other.left / 12700, other.top / 12700
                    if (shape.left / 12700 <= ox <= shape.left / 12700 + shape.width / 12700
                            and shape.top / 12700 <= oy <= shape.top / 12700 + bh):
                        tallest = max(tallest, oy + other.height / 12700
                                      - shape.top / 12700)
                if tallest > 0 and bh > tallest * T.BOX_STRETCH_MAX + T.BOX_PAD_PT:
                    errs.append("slide %d: caixa esticada (%.0fpt de altura para "
                                "%.0fpt de conteudo)" % (n, bh, tallest))
                # a caixa tambem respeita a margem: so' checar texto deixava
                # o card passar do rodape com o texto dentro dele
                bbot = (shape.top + shape.height) / 12700
                if bbot > T.CONTENT_BOTTOM_PT + 1:
                    errs.append("slide %d: caixa passa %.0fpt da margem inferior"
                                % (n, bbot - T.CONTENT_BOTTOM_PT))

            if not shape.has_text_frame:
                continue
            # Transbordo: mede PARAGRAFO A PARAGRAFO, cada um com o proprio
            # estilo. Medir a forma inteira com o estilo do primeiro run
            # superestimava grosseiramente os cards de estilos mistos.
            txt = shape.text_frame.text
            if txt.strip() and shape.width and shape.height:
                tf = shape.text_frame
                inner_w = shape.width / 12700
                inner_h = shape.height / 12700
                try:
                    inner_w -= (tf.margin_left + tf.margin_right) / 12700
                    inner_h -= (tf.margin_top + tf.margin_bottom) / 12700
                except Exception:
                    pass
                # altura = soma de (linhas x entrelinha) de cada paragrafo,
                # mais a correcao da PRIMEIRA linha do quadro quando a
                # entrelinha e' menor que o corpo (Statement usa 0,8x).
                need, known, first = 0.0, True, None
                for para in tf.paragraphs:
                    r0 = para.runs[0] if para.runs else None
                    if r0 is None or not r0.font.size or not para.line_spacing:
                        continue
                    st = by_size.get((r0.font.size.pt, bool(r0.font.bold),
                                      round(float(para.line_spacing), 3)))
                    if not st:
                        known = False
                        break
                    ptxt = "".join(r.text for r in para.runs)
                    lines = measure.wrap_lines(
                        ptxt, st["size"], bool(st["bold"]), inner_w,
                        caps=bool(st.get("caps")),
                        spacing_pt=float(st.get("spacing") or 0))
                    lead = st["size"] * st["ent"]
                    need += lines * lead
                    if first is None:
                        first = st
                    if para.space_before:
                        need += para.space_before.pt
                    if para.space_after:
                        need += para.space_after.pt
                if first is not None:
                    need += max(0.0, first["size"] - first["size"] * first["ent"])
                if known and need > inner_h + 2:
                    errs.append(
                        "slide %d: texto transborda %.0fpt (precisa %.0f, "
                        "caixa %.0f) em %r"
                        % (n, need - inner_h, need, inner_h,
                           txt[:32].replace("\n", " ")))
            # ---- nada de conteudo abaixo da margem inferior, e nada
            # encostado na tag do rodape: 2 modulos de folga.
            if txt.strip():
                bottom = (shape.top + shape.height) / 12700
                is_footer = abs(shape.top / 12700 - FOOTER_Y) < 2
                if not is_footer:
                    content_bottom = max(content_bottom, bottom)
                if bottom > T.CONTENT_BOTTOM_PT + 1:
                    errs.append("slide %d: passa %.0fpt da margem inferior: %r"
                                % (n, bottom - T.CONTENT_BOTTOM_PT,
                                   txt[:34].replace("\n", " ")))
                elif not is_footer and bottom > CONTENT_END + 1:
                    errs.append("slide %d: encosta na tag do rodape (%.0fpt de "
                                "folga, minimo %.0f): %r"
                                % (n, T.CONTENT_BOTTOM_PT - FOOTER_H - bottom,
                                   layouts.FOOTER_CLEARANCE,
                                   txt[:34].replace("\n", " ")))

            for p in shape.text_frame.paragraphs:
                if p.line_spacing is not None:
                    ls = round(float(p.line_spacing), 3)
                    if ls not in ents_ok:
                        errs.append("slide %d: entrelinha %.3f fora da tabela" % (n, ls))
                for r in p.runs:
                    if not r.text.strip():
                        continue
                    runs += 1
                    # ---- vicios de escrita (item 3): sem travessao e sem
                    # hifen duplo. Ninguem escreve com hifen duplo.
                    if "--" in r.text or "—" in r.text or "–" in r.text:
                        errs.append("slide %d: travessao ou hifen duplo em %r"
                                    % (n, r.text[:44]))
                    f = r.font
                    if f.name != T.FONT:
                        errs.append("slide %d: fonte %r em %r" % (n, f.name, r.text[:24]))
                    if f.size is None or f.size.pt not in sizes_ok:
                        errs.append("slide %d: corpo %s fora da escala em %r"
                                    % (n, f.size.pt if f.size else None, r.text[:24]))
                    try:
                        h = str(f.color.rgb).upper()
                    except Exception:
                        h = None
                    if h and h not in hexes_ok:
                        errs.append("slide %d: cor #%s fora da paleta em %r"
                                    % (n, h, r.text[:24]))

        # ---- ocupacao do canvas (item 14): conteudo amontoado no topo.
        # Aviso, nao erro: alguns arquetipos sao legitimamente curtos.
        used = (content_bottom - T.MARGIN_TOP_PT) / T.CONTENT_H_PT
        if 0 < used < T.CANVAS_FILL_MIN:
            warns.append("slide %d: conteudo ocupa so' %.0f%% da altura util. "
                         "Vazio demais." % (n, used * 100))
        elif used > T.CANVAS_FILL_MAX:
            warns.append("slide %d: conteudo ocupa %.0f%% da altura util. "
                         "Lotado: cortar texto ou dividir em dois slides."
                         % (n, used * 100))
    return len(prs.slides), runs, errs, warns


def main():
    n = generate(OUT)
    rel = os.path.relpath(OUT, REPO)
    print("== deck do design system ==")
    print("   %s" % rel)
    print("   %d slides" % n)
    print()

    slides, runs, errs, warns = validate(OUT)
    print("Validacao estrutural (tokens, transbordo, margem, caixa, escrita):")
    print("   %d slides, %d runs de texto conferidas" % (slides, runs))
    if warns:
        print()
        print("AVISOS (%d):" % len(warns))
        for w in warns[:20]:
            print("  - " + w)
        if len(warns) > 20:
            print("  ... e mais %d" % (len(warns) - 20))
    if errs:
        print()
        print("FORA DO SISTEMA (%d):" % len(errs))
        for e in errs[:40]:
            print("  - " + e)
        if len(errs) > 40:
            print("  ... e mais %d" % (len(errs) - 40))
        return 1
    print("   OK: nada fora dos tokens.")
    if slides < 50:
        print()
        print("AVISO: %d slides -- o combinado eram 50+." % slides)
    return 0


if __name__ == "__main__":
    sys.exit(main())
