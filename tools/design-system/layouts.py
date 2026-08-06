"""Construtores de slide do design system CBA 2.0.

Todo shape sai daqui, e tudo que sai daqui passa pelos tokens.

O que mudou do 1.0, tudo vindo da revisao feita a mao no PPTX:

  ritmo vertical   gaps escolhem um degrau da escala (20/40/60/80/120) e o
                   degrau tem significado. O 1.0 usava 1x em lugar de 4x.
  altura de caixa  card = padding + o MAIOR conteudo do grupo + padding.
                   Nunca esticar ate' o rodape.
  zonas            conteudo curto se distribui em zona superior + faixa
                   inferior, em vez de empilhar tudo no topo.
  lista            itens de mesma funcao vao numa caixa so', com quebra de
                   paragrafo. Uma caixa por item e' um inferno de editar.
  texto que estoura primeira frase no estilo pedido, resto um degrau abaixo.
"""

from __future__ import annotations

import copy

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

import measure
import tokens as T

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}

# atalhos da escala vertical
TIGHT, ITEM, BLOCK, SECTION, ZONE = (T.V_TIGHT, T.V_ITEM, T.V_BLOCK,
                                     T.V_SECTION, T.V_ZONE)

SATURATED = (T.AZUL, T.ROSA)   # fundos que exigem tipo branco

# O rodape (secao + numero) mora na ultima linha da area util. Nenhum conteudo
# chega perto dele: a faixa livre acima da tag e' de 3 modulos. Com 2 o texto
# de apoio ainda encostava visualmente na etiqueta.
FOOTER_H = 16
FOOTER_CLEARANCE = T.SPACING[3]
CONTENT_END = T.CONTENT_BOTTOM_PT - FOOTER_H - FOOTER_CLEARANCE

# Card da mesma cor do fundo nao e' card, e' um retangulo invisivel.
# Sobre bege o card e' branco; sobre branco, bege; sobre cor, branco.
CARD_ON = {T.BEGE: T.BRANCO, T.BRANCO: T.BEGE, T.AZUL: T.BRANCO, T.ROSA: T.BRANCO}


def card_bg_for(bg, wanted=None):
    """Cor de card valida sobre este fundo. Se pedirem a cor do proprio fundo,
    devolve a alternativa: quem quer card colorido pede azul ou rosa, e ai' a
    tipografia entra na regra de contraste sozinha."""
    if wanted and wanted.upper() != (bg or T.BEGE).upper():
        return wanted
    return CARD_ON.get((bg or T.BEGE).upper(), T.BRANCO)


def rgb(hex6: str) -> RGBColor:
    if hex6.upper() not in T.ALLOWED_HEX:
        raise ValueError("cor fora da paleta: %s" % hex6)
    return RGBColor.from_string(hex6.upper())


def style(style_id: str) -> dict:
    s = T.STYLE_BY_ID.get(style_id)
    if not s:
        raise ValueError("estilo fora da escala: %s" % style_id)
    return s


def blank_slide(prs, bg=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(bg or T.BEGE)
    return sl


def _resolve_color(st, bg_hex, override=None):
    """Cor do tipo sobre este fundo.

    Sobre fundo SATURADO (azul ou rosa) todo tipo e' branco: rosa sobre rosa
    some, e azul sobre rosa e' legivel mas nao e' o que a marca faz. So' o
    ponto final do Statement escapa, e ele tem regra propria.
    Sobre fundo claro vale o papel do estilo, com o contraste como rede.
    """
    if override:
        return override
    if bg_hex and bg_hex.upper() in SATURATED:
        return T.BRANCO
    color = T.TYPE_ROLES[st["role"]]
    if bg_hex and color.upper() == bg_hex.upper():
        return T.BRANCO
    return color


def needed(style_id, text, w):
    return measure.block_height(text, style(style_id), w)


def add_text(slide, style_id, text, x, y, w, h=None, *, bg=None, color=None,
             align="left", anchor="top", space_after=0):
    """Caixa de texto num estilo do sistema. h=None mede a altura necessaria."""
    st = style(style_id)
    if h is None:
        h = needed(style_id, text, w)
        if space_after:
            h += space_after * (text.count("\n"))
    box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = ANCHOR[anchor]
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    final_color = _resolve_color(st, bg, color)
    period_color = T.TYPE_ROLES[st["period"]] if st.get("period") is not None else None

    # O ponto final e' o unico detalhe colorido que sobrevive ao fundo
    # saturado: sobre azul ele e' rosa, sobre rosa e' azul.
    if period_color and text.rstrip().endswith("."):
        if bg and bg.upper() == T.AZUL:
            period_color = T.ROSA
        elif bg and bg.upper() == T.ROSA:
            period_color = T.AZUL
        elif bg and period_color.upper() == bg.upper():
            period_color = T.BRANCO
    else:
        period_color = None

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN[align]
        p.line_spacing = st["ent"]
        if space_after and i < len(lines) - 1:
            p.space_after = Pt(space_after)
        if period_color and i == len(lines) - 1 and line.rstrip().endswith("."):
            _run(p, line.rstrip()[:-1], st, final_color)
            _run(p, ".", st, period_color)
        else:
            _run(p, line, st, final_color)
    return box


def _run(p, text, st, color_hex):
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = T.FONT
    f.size = Pt(st["size"])
    f.bold = bool(st["bold"])
    f.color.rgb = rgb(color_hex)
    rPr = r._r.get_or_add_rPr()
    if st.get("caps"):
        rPr.set("cap", "all")
    if st.get("spacing"):
        rPr.set("spc", str(int(st["spacing"] * 100)))
    return r


def add_list(slide, style_id, items, x, y, w, *, bg=None, color=None,
             gap=None):
    """Lista numa CAIXA SO'. Itens viram paragrafos, separados por espaco de
    paragrafo. Uma caixa por item era o vicio do 1.0."""
    gap = ITEM if gap is None else gap
    text = "\n".join(items)
    return add_text(slide, style_id, text, x, y, w, bg=bg, color=color,
                    space_after=gap)


def list_height(style_id, items, w, gap=None):
    gap = ITEM if gap is None else gap
    return (sum(needed(style_id, it, w) for it in items)
            + gap * max(0, len(items) - 1))


class Stack:
    """Empilha blocos medindo cada um. O y anda sozinho."""

    def __init__(self, slide, x, y, w, bg=None):
        self.slide, self.x, self.y, self.w, self.bg = slide, x, y, w, bg
        self.top = y

    def add(self, style_id, text, gap=0, **kw):
        if not text:
            return self
        self.y += gap
        w = kw.pop("w", self.w)
        h = needed(style_id, text, w)
        add_text(self.slide, style_id, text, kw.pop("x", self.x), self.y, w, h,
                 bg=self.bg, **kw)
        self.y += h
        return self

    def fit(self, style_id, text, max_bottom, gap=0, **kw):
        """Como add(), mas encaixa na altura disponivel: primeira frase no
        estilo pedido, resto um degrau abaixo."""
        if not text:
            return self
        self.y += gap
        w = kw.pop("w", self.w)
        parts = measure.fit_text(text, style_id, w, max_bottom - self.y)
        for i, (sid, txt) in enumerate(parts):
            h = needed(sid, txt, w)
            add_text(self.slide, sid, txt, kw.get("x", self.x), self.y, w, h,
                     bg=self.bg, **{k: v for k, v in kw.items() if k != "x"})
            self.y += h
            if i < len(parts) - 1:
                self.y += TIGHT
        return self



def add_box(slide, x, y, w, h, *, fill=None, radius="std", line=None, line_w=1):
    """Retangulo do sistema. radius: 'std' (20pt), 'pill' ou 'hard'."""
    shape_type = MSO_SHAPE.RECTANGLE if radius == "hard" else MSO_SHAPE.ROUNDED_RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Pt(x), Pt(y), Pt(w), Pt(h))
    if radius != "hard":
        if radius == "pill":
            shp.adjustments[0] = 0.5
        else:
            shp.adjustments[0] = min(0.5, (T.RADIUS_PT / (min(w, h) / 2)) * 0.5)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False          # anti-vicios: nada de sombra
    shp.text_frame.word_wrap = True
    return shp


def _kicker_style(text):
    """Numero de item e' informacao, nao etiqueta: entra em Titulo 60.
    Rotulo textual (VARIACAO, PAL0) continua CAPS 12."""
    t = (text or "").strip().rstrip(".")
    return "dsH1" if t.isdigit() else "dsCaps12"


def card_blocks(item, *, kicker_style=None):
    """(estilo, texto, gap_antes) de um item de card.

    O gap entre o titulo e o texto de apoio e' TIGHT: eles sao o mesmo bloco
    de leitura. Gap maior desmancha o par.
    """
    out = []
    if item.get("kicker"):
        out.append((kicker_style or _kicker_style(item["kicker"]), item["kicker"], 0))
    if item.get("title"):
        out.append(("dsH4", item["title"], ITEM if out else 0))
    if item.get("body"):
        out.append(("dsCorpoPilar", item["body"], TIGHT if out else 0))
    return out


def blocks_height(blocks, w):
    h = 0.0
    for i, (sid, txt, gap) in enumerate(blocks):
        if i:
            h += gap
        h += needed(sid, txt, w)
    return h


def add_card(slide, x, y, w, h, blocks, *, fill=None, radius="std", line=None,
             pad=None):
    """Card com o texto DENTRO da forma.

    Caixa de texto sobreposta ao retangulo obriga quem edita a mexer em dois
    objetos e a manter os dois alinhados na mao. Uma forma so', com paragrafos
    de estilos diferentes, edita como um bloco.
    """
    pad = T.BOX_PAD_PT if pad is None else pad
    shp = add_box(slide, x, y, w, h, fill=fill, radius=radius, line=line)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(pad)
    tf.margin_top = tf.margin_bottom = Pt(pad)
    tf.vertical_anchor = ANCHOR["top"]
    for i, (sid, txt, gap) in enumerate(blocks):
        st = style(sid)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN["left"]
        p.line_spacing = st["ent"]
        if i and gap:
            p.space_before = Pt(gap)
        _run(p, txt, st, _resolve_color(st, fill))
    return shp


def footer(slide, section, number, bg=None):
    """Rodape de localizacao. Fica na ultima linha da area util."""
    y = T.CONTENT_BOTTOM_PT - 16
    if section:
        add_text(slide, "dsCaps12", section, T.MARGIN_LEFT_PT, y, 700, 16, bg=bg)
    add_text(slide, "dsLegenda12", str(number), T.MARGIN_RIGHT_PT - 100, y, 100,
             16, bg=bg, align="right")


def place_group(top, bottom, group_h):
    """Onde comeca um grupo de altura group_h no espaco de top a bottom.

    Sem esticar. Se sobra espaco, o grupo CENTRALIZA no que sobrou; se nao
    sobra, comeca no topo. Esticar o gap para encher a pagina foi o erro do
    2.0: virava 120pt entre linhas de uma grade cujo gutter e' 40.
    """
    room = bottom - top - group_h
    if room <= 0:
        return top
    return top + room / 2


def grid_height(rows, block_h, gap):
    return rows * block_h + (rows - 1) * gap


# ---------------------------------------------------------------- arquetipos

def hero_cover(prs, spec):
    """Capa. Fundo SEMPRE saturado: azul (preferencial) ou rosa."""
    bg = spec.get("bg", T.AZUL)
    if bg not in SATURATED:
        raise ValueError("capa exige fundo azul ou rosa, recebeu %s" % bg)
    sl = blank_slide(prs, bg)
    w = T.col_span(3)
    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, w, bg)
    s.add("dsLegenda12", spec.get("eyebrow"))
    s.add(spec.get("style", "dsHero"), spec["title"], gap=SECTION)
    # linha de apoio da capa e' curta e de destaque: Destaque 44 e' o papel dele
    s.add("dsCorpo", spec.get("sub"), gap=SECTION, w=T.col_span(2))
    return sl


def chapter_divider(prs, spec):
    """Abertura de bloco. Numero em '1.', nunca '01'."""
    bg = spec.get("bg", T.AZUL)
    sl = blank_slide(prs, bg)
    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, T.col_span(3), bg)
    s.add("dsBigNumber", spec.get("number"))
    s.add("dsMega", spec["title"], gap=SECTION)
    s.add("dsTexto34", spec.get("sub"), gap=SECTION, w=T.col_span(2))
    return sl


def spec_page(prs, spec):
    """Titulo + intro + pares rotulo/valor, distribuidos ate' a base."""
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, T.CONTENT_W_PT, bg)
    s.add("dsCaps12", spec.get("eyebrow"))
    s.add("dsH1", spec["title"], gap=ITEM if spec.get("eyebrow") else 0)
    if spec.get("intro"):
        s.add("dsTexto34", spec["intro"], gap=BLOCK, w=T.col_span(3))

    rows = spec.get("rows", [])
    if rows:
        # ate' 8 pares cabem em 2 blocos de 2 colunas; acima disso, 4 blocos
        # de 1 coluna. Apertar a linha em vez de abrir coluna era o que fazia
        # a tabela de entrelinha estourar a pagina.
        ncols = 2 if len(rows) <= 8 else 4
        span = T.COLUMNS // ncols
        col_w = T.col_span(span)
        per_col = (len(rows) + ncols - 1) // ncols
        lab_h = max(needed("dsCaps12", k, col_w) for k, _ in rows)
        val_h = max(needed("dsH4", v, col_w) for _, v in rows)
        row_h = lab_h + TIGHT + val_h
        # distribui as linhas ate' a base da area util
        top = s.y + BLOCK
        bottom = CONTENT_END
        pitch = row_h + ITEM
        top = place_group(top, bottom, grid_height(per_col, row_h, ITEM))
        for i, (label, value) in enumerate(rows):
            col, row = divmod(i, per_col)
            rx = T.col_x(col * span)
            ry = top + row * pitch
            add_text(sl, "dsCaps12", label, rx, ry, col_w, bg=bg)
            add_text(sl, "dsH4", value, rx, ry + lab_h + TIGHT, col_w, bg=bg)
    return sl


def type_specimen(prs, spec):
    """Um estilo por slide: amostra no proprio corpo + ficha tecnica."""
    st = style(spec["style_id"])
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)

    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, T.CONTENT_W_PT, bg)
    s.add("dsCaps12", spec.get("eyebrow", "ESCALA TIPOGRÁFICA"))
    s.add("dsH1", st["label"], gap=ITEM)

    cells = [("CORPO", "%gpt" % st["size"]),
             ("PESO", "Bold" if st["bold"] else "Regular"),
             ("COR", T.PALETTE[0 if st["role"] == 0 else 1]["name"].capitalize()),
             ("ENTRELINHA", "%.2fx" % st["ent"]),
             ("PX @1080", "%dpx" % round(T.pt_to_px(st["size"])))]
    use_w = T.col_span(3)
    use_h = needed("dsCorpoPilar", st["use"], use_w)
    band_h = 12 + TIGHT + 28
    fy = CONTENT_END - use_h - BLOCK - band_h
    cw = T.CONTENT_W_PT / len(cells)
    for i, (k, v) in enumerate(cells):
        cx = T.MARGIN_LEFT_PT + i * cw
        add_text(sl, "dsCaps12", k, cx, fy, cw - ITEM, bg=bg)
        add_text(sl, "dsH4", v, cx, fy + 12 + TIGHT, cw - ITEM, bg=bg)
    add_text(sl, "dsCorpoPilar", st["use"], T.MARGIN_LEFT_PT,
             fy + band_h + BLOCK, use_w, use_h, bg=bg)

    top = s.y + BLOCK
    avail = fy - BLOCK - top
    sh = needed(spec["style_id"], spec["sample"], T.CONTENT_W_PT)
    add_text(sl, spec["style_id"], spec["sample"], T.MARGIN_LEFT_PT,
             top + max(0, (avail - sh) / 2), T.CONTENT_W_PT, sh, bg=bg)
    return sl


def swatch_page(prs, spec):
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    pal = spec["palette"]
    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, T.CONTENT_W_PT, bg)
    s.add("dsCaps12", "PALETA")
    s.add("dsH1", pal["name"].capitalize(), gap=ITEM)

    sw_y = s.y + BLOCK
    sw_w = T.col_span(2)
    sw_h = CONTENT_END - sw_y
    # contorno SO' quando a amostra e' a propria cor do fundo. Branco sobre
    # bege ja' se distingue e nao precisa de borda.
    same = pal["hex"].upper() == (bg or T.BEGE).upper()
    add_box(sl, T.MARGIN_LEFT_PT, sw_y, sw_w, sw_h, fill=pal["hex"],
            line=T.AZUL if same else None)

    tx = T.col_x(2)
    tw = T.col_span(2)
    ts = Stack(sl, tx, sw_y, tw, bg)
    for i, (k, v) in enumerate([("HEX", "#" + pal["hex"]),
                                ("RGB", "%d, %d, %d" % pal["rgb"]),
                                ("SLOT", "pal%d" % spec["index"])]):
        ts.add("dsCaps12", k, gap=BLOCK if i else 0)
        ts.add("dsH4", v, gap=TIGHT)
    ts.add("dsCorpoPilar", pal["role"], gap=SECTION)
    return sl


def multi_card_grid(prs, spec):
    """2 a 4 cards. Altura pelo MAIOR conteudo, nunca ate' o rodape."""
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, T.CONTENT_W_PT, bg)
    s.add("dsCaps12", spec.get("eyebrow"))
    s.add("dsH1", spec["title"], gap=ITEM)

    cards = spec["cards"]
    n = len(cards)
    span = T.COLUMNS // n if T.COLUMNS % n == 0 else 1
    cw = T.col_span(span) if T.COLUMNS % n == 0 else \
        (T.CONTENT_W_PT - T.GUTTER_PT * (n - 1)) / n
    pad = T.BOX_PAD_PT
    inner = cw - pad * 2
    blocks = [card_blocks(c) for c in cards]
    ch = max(blocks_height(b, inner) for b in blocks) + pad * 2
    cy = s.y + BLOCK
    card_bg = card_bg_for(bg, spec.get("card_bg"))
    for i, b in enumerate(blocks):
        cx = (T.col_x(i * span) if T.COLUMNS % n == 0
              else T.MARGIN_LEFT_PT + i * (cw + T.GUTTER_PT))
        add_card(sl, cx, cy, cw, ch, b, fill=card_bg, pad=pad)
    return sl


def card_grid_5(prs, spec):
    """Grade densa de itens curtos, em cards."""
    return _grid(prs, spec, cards=True)


def grid_plain(prs, spec):
    """A mesma grade, SEM caixa. Alternativa mais limpa quando o item e' curto."""
    return _grid(prs, spec, cards=False)


def _grid(prs, spec, cards):
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, T.CONTENT_W_PT, bg)
    s.add("dsH1", spec["title"])
    s.add("dsCorpoPilar", spec.get("intro"), gap=ITEM, w=T.col_span(3))

    items = spec["items"]
    cols = spec.get("cols", 4)
    span = max(1, T.COLUMNS // cols)
    cw = T.col_span(span) if cols <= T.COLUMNS else \
        (T.CONTENT_W_PT - T.GUTTER_PT * (cols - 1)) / cols
    pad = T.BOX_PAD_PT if cards else 0
    rows = (len(items) + cols - 1) // cols
    top = s.y + BLOCK
    bottom = CONTENT_END

    # Grade densa e' conteudo MICRO: se o corpo em Apoio 20 nao couber,
    # desce para Texto 15, que existe pra isso (item 4 da revisao).
    inner = cw - pad * 2
    for candidate in ("dsCorpoPilar", "dsTexto15"):
        blocks = [card_blocks(it) for it in items]
        for b in blocks:
            for j, (sid, txt, gap) in enumerate(b):
                if sid == "dsCorpoPilar":
                    b[j] = (candidate, txt, gap)
        ch = max(blocks_height(b, inner) for b in blocks) + pad * 2
        if top + rows * ch + (rows - 1) * T.GUTTER_PT <= bottom:
            break
    # o vao vertical e' o MESMO do horizontal. Grade com gutter 40 na
    # horizontal e 120 na vertical nao le' como grade.
    pitch = ch + T.GUTTER_PT
    top = place_group(top, bottom, grid_height(rows, ch, T.GUTTER_PT))
    card_bg = card_bg_for(bg, spec.get("card_bg"))
    for i, _it in enumerate(items):
        r, c = divmod(i, cols)
        cx = (T.col_x(c * span) if cols <= T.COLUMNS
              else T.MARGIN_LEFT_PT + c * (cw + T.GUTTER_PT))
        cy = top + r * pitch
        if cards:
            add_card(sl, cx, cy, cw, ch, blocks[i], fill=card_bg, pad=pad)
        else:
            cs = Stack(sl, cx, cy, cw, bg)
            for j, (sid, txt, gap) in enumerate(blocks[i]):
                cs.add(sid, txt, gap=gap if j else 0)
    return sl


def zoned_content(prs, spec):
    """Zona superior (argumento) + faixa inferior (desdobramento).

    Substitui o pillar_card_dense do 1.0, que empilhava tudo no topo e deixava
    metade do slide vazia.
    """
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    items = spec.get("items", [])
    n = max(1, len(items))
    span = max(1, T.COLUMNS // n)
    iw = T.col_span(span)

    # Mede a zona de cima e a faixa, e posiciona as duas como UM grupo.
    # Ancorar a faixa na base jogava os blocos para o pe' da pagina.
    tw = T.col_span(3)
    top_blocks = []
    if spec.get("eyebrow"):
        top_blocks.append(("dsEyebrow", spec["eyebrow"], 0))
    top_blocks.append(("dsH1", spec["title"], ITEM if top_blocks else 0))
    if spec.get("body"):
        top_blocks.append(("dsTexto34", spec["body"], BLOCK))
    top_h = blocks_height(top_blocks, tw)

    band, band_h = None, 0
    for candidate in ("dsCorpoPilar", "dsTexto15"):
        band = [card_blocks(it, kicker_style="dsCaps12") for it in items]
        for b in band:
            for j, (sid, txt, gap) in enumerate(b):
                if sid == "dsCorpoPilar":
                    b[j] = (candidate, txt, gap)
        band_h = max(blocks_height(b, iw) for b in band) if band else 0
        if T.MARGIN_TOP_PT + top_h + BLOCK + band_h <= CONTENT_END:
            break

    total = top_h + BLOCK + band_h
    y0 = place_group(T.MARGIN_TOP_PT, CONTENT_END, total)

    st = Stack(sl, T.MARGIN_LEFT_PT, y0, tw, bg)
    for i, (sid, txt, gap) in enumerate(top_blocks):
        st.add(sid, txt, gap=gap if i else 0)
    by = y0 + top_h + BLOCK
    for i, b in enumerate(band):
        cs = Stack(sl, T.col_x(i * span), by, iw, bg)
        for j, (sid, txt, gap) in enumerate(b):
            cs.add(sid, txt, gap=gap if j else 0)
    return sl


def stat_band(prs, spec):
    """Zona superior + faixa inferior de NUMEROS. Variacao do zoned_content."""
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    stats = spec["stats"]
    n = len(stats)
    span = max(1, T.COLUMNS // n)
    iw = T.col_span(span)

    tw = T.col_span(3)
    top_blocks = []
    if spec.get("eyebrow"):
        top_blocks.append(("dsCaps12", spec["eyebrow"], 0))
    top_blocks.append(("dsH1", spec["title"], ITEM if top_blocks else 0))
    if spec.get("body"):
        top_blocks.append(("dsTexto34", spec["body"], BLOCK))
    top_h = blocks_height(top_blocks, tw)

    band = [[("dsMega", x["value"], 0), ("dsCorpoPilar", x["label"], TIGHT)]
            for x in stats]
    band_h = max(blocks_height(b, iw) for b in band)

    y0 = place_group(T.MARGIN_TOP_PT, CONTENT_END, top_h + BLOCK + band_h)
    st = Stack(sl, T.MARGIN_LEFT_PT, y0, tw, bg)
    for i, (sid, txt, gap) in enumerate(top_blocks):
        st.add(sid, txt, gap=gap if i else 0)
    by = y0 + top_h + BLOCK
    for i, b in enumerate(band):
        cs = Stack(sl, T.col_x(i * span), by, iw, bg)
        for j, (sid, txt, gap) in enumerate(b):
            cs.add(sid, txt, gap=gap if j else 0)
    return sl


def quote_side_image(prs, spec):
    """Texto a' esquerda, painel de cor a' direita ocupando a altura util."""
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    tw = T.col_span(2)
    px = T.col_x(2)
    pw = T.col_span(2)
    ph = T.CONTENT_H_PT
    panel_bg = card_bg_for(bg, spec.get("panel_bg", T.AZUL))
    add_box(sl, px, T.MARGIN_TOP_PT, pw, ph, fill=panel_bg)
    if spec.get("panel_text"):
        pstyle = spec.get("panel_style", "dsH3")
        pad = T.BOX_PAD_PT
        th = needed(pstyle, spec["panel_text"], pw - pad * 2)
        add_text(sl, pstyle, spec["panel_text"], px + pad,
                 T.MARGIN_TOP_PT + max(pad, (ph - th) / 2), pw - pad * 2, th,
                 bg=panel_bg)

    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, tw, bg)
    s.add("dsEyebrow", spec.get("eyebrow"))
    s.add("dsH1", spec["title"], gap=ITEM)
    s.add("dsTexto34", spec["body"], gap=BLOCK)
    return sl


def do_dont(prs, spec):
    """Duas colunas. Cabecalho em Subtitulo 34 e itens numa CAIXA SO'."""
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, T.CONTENT_W_PT, bg)
    s.add("dsH1", spec["title"])
    s.add("dsCorpoPilar", spec.get("intro"), gap=ITEM, w=T.col_span(3))

    cw = T.col_span(2)
    pad = T.BOX_PAD_PT
    inner = cw - pad * 2
    cols = [(spec.get("do_title", "Sim"), spec["do"], T.AZUL),
            (spec.get("dont_title", "Nunca"), spec["dont"], T.ROSA)]
    card_bg = card_bg_for(bg, spec.get("card_bg"))
    ch = pad * 2 + max(needed("dsH3", h, inner) + BLOCK
                       + list_height("dsCorpoPilar", it, inner)
                       for h, it, _ in cols)
    cy = s.y + BLOCK
    for i, (head, items, tone) in enumerate(cols):
        cx = T.col_x(i * 2)
        shp = add_card(sl, cx, cy, cw, ch,
                       [("dsH3", head, 0)], fill=card_bg, pad=pad)
        # o cabecalho vai colorido; a lista entra no mesmo objeto
        shp.text_frame.paragraphs[0].runs[0].font.color.rgb = rgb(tone)
        for j, it in enumerate(items):
            p = shp.text_frame.add_paragraph()
            p.line_spacing = style("dsCorpoPilar")["ent"]
            p.space_before = Pt(BLOCK if j == 0 else ITEM)
            _run(p, it, style("dsCorpoPilar"),
                 _resolve_color(style("dsCorpoPilar"), card_bg))
    return sl


def diagram_page(prs, spec):
    """Esquema em escala: os retangulos sao proporcionais de verdade."""
    bg = spec.get("bg", T.BEGE)
    sl = blank_slide(prs, bg)
    lw = T.col_span(2)     # 1 coluna quebrava titulo de 60pt no meio da palavra
    s = Stack(sl, T.MARGIN_LEFT_PT, T.MARGIN_TOP_PT, lw, bg)
    s.add("dsH1", spec["title"])
    s.add("dsCorpoPilar", spec.get("intro", ""), gap=BLOCK)

    legend = spec.get("legend", [])
    dw_guess = T.col_span(2)
    leg_blocks = [[("dsCaps12", k, 0), ("dsCorpoPilar", v, TIGHT)]
                  for k, v in legend]
    leg_h = (sum(blocks_height(b, dw_guess) for b in leg_blocks)
             + ITEM * max(0, len(legend) - 1)) if legend else 0
    dy = T.MARGIN_TOP_PT + BLOCK
    max_dh = CONTENT_END - dy - (leg_h + BLOCK if legend else 0)
    dw = min(T.col_span(2), max_dh * T.PAGE_W_PT / T.PAGE_H_PT)
    dh = dw * T.PAGE_H_PT / T.PAGE_W_PT
    dx = T.MARGIN_RIGHT_PT - dw
    scale = dw / T.PAGE_W_PT
    add_box(sl, dx, dy, dw, dh, fill=T.BRANCO, radius="hard", line=T.AZUL)
    if spec.get("show_margin", True):
        m = T.MARGIN_PT * scale
        add_box(sl, dx + m, dy + m, dw - 2 * m, dh - 2 * m,
                radius="hard", line=T.ROSA)
    if spec.get("show_module"):
        u = T.SPACING_UNIT_PT * scale
        for i in range(8):
            add_box(sl, dx + T.MARGIN_PT * scale + i * u * 2, dy + dh - u * 3,
                    u, u, fill=T.AZUL, radius="hard")

    if legend:
        ls = Stack(sl, dx, min(dy + dh + BLOCK, CONTENT_END - leg_h), dw, bg)
        for i, (label, value) in enumerate(legend):
            ls.add("dsCaps12", label, gap=ITEM if i else 0)
            ls.add("dsCorpoPilar", value, gap=TIGHT)
    return sl


def closing(prs, spec):
    """Fecho. Sem linha de versao: o 'ultimo slide' de verdade vira arquetipo
    proprio mais adiante."""
    bg = spec.get("bg", T.AZUL)
    sl = blank_slide(prs, bg)
    w = T.col_span(3)
    th = needed("dsHero", spec["title"], w)
    add_text(sl, "dsHero", spec["title"], T.MARGIN_LEFT_PT,
             (T.PAGE_H_PT - th) / 2, w, th, bg=bg)
    return sl


BUILDERS = {
    "hero_cover": hero_cover,
    "chapter_divider": chapter_divider,
    "spec_page": spec_page,
    "type_specimen": type_specimen,
    "swatch_page": swatch_page,
    "multi_card_grid": multi_card_grid,
    "card_grid_5": card_grid_5,
    "grid_plain": grid_plain,
    "zoned_content": zoned_content,
    "stat_band": stat_band,
    "quote_side_image": quote_side_image,
    "do_dont": do_dont,
    "diagram_page": diagram_page,
    "closing": closing,
}


def build(prs, spec, index):
    kind = spec["kind"]
    fn = BUILDERS.get(kind)
    if not fn:
        raise ValueError("arquetipo desconhecido: %s" % kind)
    sl = fn(prs, copy.deepcopy(spec))
    if spec.get("footer", True) and kind not in ("hero_cover", "closing"):
        footer(sl, spec.get("section", ""), index, bg=spec.get("bg", T.BEGE))
    return sl
