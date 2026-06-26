#!/usr/bin/env python3
"""
Build assets/estilos-bg.pptx — the "style anchor" deck for the Typography module.

One slide holds the B+G typographic styles as REAL, exactly-formatted text boxes
(font, weight, size, color, the colored final period AND the line spacing baked in,
as a MULTIPLE — spcPct). The add-in inserts this slide via
insertSlidesFromBase64(KeepSourceFormatting); the user then copy/pastes a box or
uses the Format Painter onto their object — both native paths carry the line
spacing that the Office.js API cannot set.

Run:  python3 tools/build-styles-reference.py
Out:  assets/estilos-bg.pptx
"""
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

RED = RGBColor(0xFC, 0x5E, 0x6D)      # cor 1
BLUE = RGBColor(0x43, 0x6A, 0xE1)     # cor 2
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x76, 0x76, 0x7C)
FONT = "Avenir Next"

# col: 0 = left, 1 = right
# (key, label, sample, size, line_spacing, bold, text_color, period_color, fill, col)
STYLES = [
    ("hero",       "HERO STATEMENT · 120pt · 0.8x",        "Ideia grande",      120, 0.80, True,  RED,   BLUE, None, 0),
    ("hero_blue",  "HERO sobre 436AE1 · branco",           "Ideia grande",      120, 0.80, True,  WHITE, RED,  BLUE, 0),
    ("mega",       "MEGA STATEMENT · 80pt · 0.9x",         "Mega statement",     80, 0.90, True,  BLUE,  None, None, 0),
    ("h1",         "H1 TÍTULO DE PÁGINA · 60pt · 0.9x",    "Título de página",   60, 0.90, True,  RED,   None, None, 0),
    ("label_sec",  "LABEL DE SEÇÃO · 60pt · 0.9x",         "Label de seção",     60, 0.90, True,  RED,   None, None, 1),
    ("corpo",      "CORPO DE TEXTO · 44pt · 1.15x · Reg",  "Corpo de texto",     44, 1.15, False, BLUE,  None, None, 1),
    ("h3",         "H3 CORPO DESCRITIVO · 34pt · 0.95x",   "Corpo descritivo",   34, 0.95, True,  RED,   None, None, 1),
    ("h4",         "H4 SUBTÍTULO DE PILAR · 28pt · 1.0x",  "Subtítulo de pilar", 28, 1.00, True,  BLUE,  None, None, 1),
    ("h5",         "H5 TEXTO DESCRITIVO · 24pt · 1.0x · Reg", "Texto descritivo", 24, 1.00, False, BLUE, None, None, 1),
    ("corpo_pil",  "CORPO DE PILAR · 20pt · 1.3x · Reg",   "Corpo de pilar",     20, 1.30, False, BLUE,  None, None, 1),
    ("eyebrow",    "LABEL / EYEBROW · 18pt · 1.0x",        "Label / eyebrow",    18, 1.00, True,  RED,   None, None, 1),
    ("caption",    "LEGENDA / CAPTION · 16pt · 1.3x · Reg","Legenda / caption",  16, 1.30, False, BLUE,  None, None, 1),
]

# Slide ~ CBA B+G page (16:9, large) so the 120pt boxes don't collide.
SLIDE_W, SLIDE_H = 1583, 890
COL_X = {0: 70, 1: 812}
COL_W = 700
CAP_H = 22
GAP = 16


def add_caption(slide, x, top, text):
    box = slide.shapes.add_textbox(Pt(x), Pt(top), Pt(COL_W), Pt(CAP_H))
    p = box.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    f = r.font; f.name = FONT; f.size = Pt(12); f.bold = True; f.color.rgb = GRAY


def add_style_box(slide, x, top, height, sample, size, spacing, bold, tcol, pcol, fill):
    box = slide.shapes.add_textbox(Pt(x), Pt(top), Pt(COL_W), Pt(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)

    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    else:
        box.fill.background(); box.line.fill.background()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = spacing          # float => MULTIPLE (spcPct), e.g. 0.8 == 80%
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    r = p.add_run(); r.text = sample
    f = r.font; f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = tcol

    if pcol is not None:
        d = p.add_run(); d.text = "."
        df = d.font; df.name = FONT; df.size = Pt(size); df.bold = bold; df.color.rgb = pcol


def main():
    prs = Presentation()
    prs.slide_width = Pt(SLIDE_W)
    prs.slide_height = Pt(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    head = slide.shapes.add_textbox(Pt(70), Pt(24), Pt(1450), Pt(30))
    hr = head.text_frame.paragraphs[0].add_run()
    hr.text = "Estilos B+G — copie um bloco ou use o Pincel de Formatação (carrega a entrelinha). Apague este slide depois."
    hf = hr.font; hf.name = FONT; hf.size = Pt(15); hf.bold = True; hf.color.rgb = GRAY

    y = {0: 72, 1: 72}
    for key, label, sample, size, spacing, bold, tcol, pcol, fill, col in STYLES:
        box_h = round(size * 1.32) + 6
        x = COL_X[col]
        add_caption(slide, x, y[col], label)
        add_style_box(slide, x, y[col] + CAP_H, box_h, sample, size, spacing, bold, tcol, pcol, fill)
        y[col] += CAP_H + box_h + GAP

    out = os.path.normpath(os.path.join(os.path.dirname(__file__), "..", "assets", "estilos-bg.pptx"))
    prs.save(out)
    print("wrote", out, "| col0 bottom", y[0], "col1 bottom", y[1], "(slide", SLIDE_H, "pt)")


if __name__ == "__main__":
    main()
